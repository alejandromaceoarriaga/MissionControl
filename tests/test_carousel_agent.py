# tests/test_carousel_agent.py
import os
from agent.sub_agents.carousel_agent import CarouselAgent
from pptx import Presentation

SAMPLE_DATA = {
    "bcrd": {
        "tpm": {"value": 5.25},
        "tasas_bancarias": {
            "bancos_multiples": {"activa": 13.27, "pasiva": 6.14},
            "aayp": {"activa": 14.5, "pasiva": 5.8},
            "bancos_ahorro_credito": {"activa": 15.0, "pasiva": 5.5},
            "interbancaria": 6.06,
        },
        "imae": {"var_interanual": 5.3},
        "inflacion": {"var_interanual": 3.8},
        "tipo_cambio": {"compra": 60.55, "venta": 61.19},
        "reservas": {"brutas_mm_usd": 16143.1},
    },
    "sb": {
        "rankings": {
            "top_cartera_growth": [
                {"entidad": "BanReservas", "_growth_pct": 18.5},
                {"entidad": "Popular", "_growth_pct": 14.2},
                {"entidad": "BHD", "_growth_pct": 11.8},
            ],
            "top_npl_increase": [
                {"entidad": "Banco X", "_growth_pct": 0.9},
                {"entidad": "Banco Y", "_growth_pct": 0.5},
            ],
        }
    },
    "imf": {"capital_adequacy": {"value": 17.5, "year": "2024"}},
    "wb": {"credit_gdp": [{"year": "2024", "value": 31.92}]},
}

def test_generate_creates_pptx(tmp_path):
    agent = CarouselAgent(output_dir=str(tmp_path))
    result = agent.generate(SAMPLE_DATA, period="Abril 2026")
    assert result.endswith(".pptx")
    assert os.path.exists(result)
    assert os.path.getsize(result) > 5000

def test_generate_has_correct_slide_count(tmp_path):
    agent = CarouselAgent(output_dir=str(tmp_path))
    result = agent.generate(SAMPLE_DATA, period="Abril 2026")
    prs = Presentation(result)
    assert len(prs.slides) >= 8

def test_slide_dimensions_are_square(tmp_path):
    agent = CarouselAgent(output_dir=str(tmp_path))
    result = agent.generate(SAMPLE_DATA, period="Abril 2026")
    prs = Presentation(result)
    # Square format for LinkedIn carousel (width == height)
    assert prs.slide_width == prs.slide_height
