# agent/sub_agents/carousel_agent.py
import os
from datetime import date
from io import BytesIO

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors
DARK_BLUE = RGBColor(0x1a, 0x1a, 0x2e)
MED_BLUE = RGBColor(0x31, 0x82, 0xce)
LIGHT_BLUE = RGBColor(0x63, 0xb3, 0xed)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x38, 0xa1, 0x69)
RED = RGBColor(0xe5, 0x3e, 0x3e)
YELLOW = RGBColor(0xd6, 0x9e, 0x2e)
GRAY = RGBColor(0x71, 0x80, 0x96)
LIGHT_GRAY = RGBColor(0xf7, 0xfa, 0xfc)

# 1080x1080px at 96dpi ~= 11.25 inches
SLIDE_SIZE = Inches(11.25)


class CarouselAgent:
    def __init__(self, output_dir: str = "outputs/carousel"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def _new_prs(self) -> Presentation:
        prs = Presentation()
        prs.slide_width = SLIDE_SIZE
        prs.slide_height = SLIDE_SIZE
        return prs

    def _blank_slide(self, prs: Presentation, bg_color: RGBColor = None):
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        if bg_color:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
        return slide

    def _txt(self, slide, text: str, left, top, width, height,
             size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def _chart_png(self, fig) -> BytesIO:
        buf = BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='#f7fafc')
        buf.seek(0)
        plt.close(fig)
        return buf

    def _slide_cover(self, prs, period):
        slide = self._blank_slide(prs, DARK_BLUE)
        m = Inches(0.5)
        w = SLIDE_SIZE - 2 * m
        self._txt(slide, "SISTEMA BANCARIO DOMINICANO", m, Inches(2.5), w, Inches(0.7),
                  size=18, bold=True, color=MED_BLUE, align=PP_ALIGN.CENTER)
        self._txt(slide, "Barometro Mensual", m, Inches(3.5), w, Inches(1.5),
                  size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        self._txt(slide, period, m, Inches(5.2), w, Inches(0.9),
                  size=26, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
        self._txt(slide, "Fuente: BCRD · SB · IMF · World Bank", m, Inches(10.2), w, Inches(0.5),
                  size=11, color=GRAY, align=PP_ALIGN.CENTER)

    def _slide_kpis(self, prs, data):
        slide = self._blank_slide(prs, WHITE)
        m = Inches(0.3)
        w = SLIDE_SIZE - 2 * m
        self._txt(slide, "KPIs CLAVE DEL MES", m, Inches(0.3), w, Inches(0.7),
                  size=22, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        bcrd = data.get("bcrd", {})
        tasas = bcrd.get("tasas_bancarias", {}).get("bancos_multiples", {})
        kpis = [
            ("TPM", f"{bcrd.get('tpm', {}).get('value', 'N/D')}%"),
            ("Activa BM", f"{tasas.get('activa', 'N/D')}%"),
            ("Pasiva BM", f"{tasas.get('pasiva', 'N/D')}%"),
            ("IMAE i.a.", f"+{bcrd.get('imae', {}).get('var_interanual', 'N/D')}%"),
            ("Inflacion i.a.", f"{bcrd.get('inflacion', {}).get('var_interanual', 'N/D')}%"),
            ("USD/DOP", f"RD${bcrd.get('tipo_cambio', {}).get('venta', 'N/D')}"),
        ]
        card_w, card_h = Inches(3.4), Inches(3.8)
        gap = Inches(0.15)
        start_x, start_y = Inches(0.35), Inches(1.3)
        for i, (label, value) in enumerate(kpis):
            col, row = i % 3, i // 3
            x = start_x + col * (card_w + gap)
            y = start_y + row * (card_h + gap)
            shape = slide.shapes.add_shape(1, x, y, card_w, card_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_GRAY
            shape.line.color.rgb = MED_BLUE
            shape.line.width = Pt(2)
            self._txt(slide, label, x + Inches(0.15), y + Inches(0.2),
                      card_w - Inches(0.3), Inches(0.6), size=13, color=GRAY)
            self._txt(slide, value, x + Inches(0.1), y + Inches(0.9),
                      card_w - Inches(0.2), Inches(1.5), size=34, bold=True, color=DARK_BLUE)

    def _slide_macro_chart(self, prs, data):
        slide = self._blank_slide(prs, WHITE)
        m = Inches(0.3)
        w = SLIDE_SIZE - 2 * m
        self._txt(slide, "ENTORNO MACROECONÓMICO", m, Inches(0.3), w, Inches(0.6),
                  size=20, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        bcrd = data.get("bcrd", {})
        tasas = bcrd.get("tasas_bancarias", {}).get("bancos_multiples", {})
        labels = ['TPM', 'Activa BM', 'Pasiva BM', 'IMAE i.a.', 'Inflacion']
        values = [
            bcrd.get("tpm", {}).get("value") or 0,
            tasas.get("activa") or 0,
            tasas.get("pasiva") or 0,
            bcrd.get("imae", {}).get("var_interanual") or 0,
            bcrd.get("inflacion", {}).get("var_interanual") or 0,
        ]
        colors = ['#1a1a2e', '#3182ce', '#63b3ed', '#38a169', '#d69e2e']
        fig, ax = plt.subplots(figsize=(9, 6.5))
        bars = ax.bar(labels, values, color=colors, width=0.6, edgecolor='white')
        for bar, val in zip(bars, values):
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.08,
                    f'{val:.1f}%', ha='center', va='bottom', fontsize=11, fontweight='bold')
        ax.set_ylabel('%', fontsize=12)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        buf = self._chart_png(fig)
        slide.shapes.add_picture(buf, Inches(0.4), Inches(1.2), Inches(10.4), Inches(7.5))
        tc = bcrd.get("tipo_cambio", {})
        res = bcrd.get("reservas", {})
        self._txt(slide, f"USD/DOP: RD${tc.get('venta','N/D')} (venta)  ·  Reservas: US${res.get('brutas_mm_usd','N/D')}MM",
                  m, Inches(10.2), w, Inches(0.5), size=12, color=GRAY, align=PP_ALIGN.CENTER)

    def _slide_tasas(self, prs, data):
        slide = self._blank_slide(prs, WHITE)
        m = Inches(0.3)
        w = SLIDE_SIZE - 2 * m
        self._txt(slide, "TASAS DE INTERÉS POR TIPO DE ENTIDAD", m, Inches(0.3), w, Inches(0.6),
                  size=18, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        bcrd = data.get("bcrd", {})
        tasas = bcrd.get("tasas_bancarias", {})
        bm = tasas.get("bancos_multiples", {})
        aayp = tasas.get("aayp", {})
        bac = tasas.get("bancos_ahorro_credito", {})
        categories = ['Bancos\nMultiples', 'AAyP', 'Bancos Ahorro\ny Credito']
        activas = [bm.get("activa") or 0, aayp.get("activa") or 0, bac.get("activa") or 0]
        pasivas = [bm.get("pasiva") or 0, aayp.get("pasiva") or 0, bac.get("pasiva") or 0]
        x = range(len(categories))
        fig, ax = plt.subplots(figsize=(9, 6.5))
        bars1 = ax.bar([xi - 0.2 for xi in x], activas, 0.35, label='Activa', color='#1a1a2e')
        bars2 = ax.bar([xi + 0.2 for xi in x], pasivas, 0.35, label='Pasiva', color='#3182ce')
        ax.set_xticks(list(x))
        ax.set_xticklabels(categories, fontsize=11)
        ax.set_ylabel('%', fontsize=12)
        ax.legend(fontsize=11)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        for bars in [bars1, bars2]:
            for bar in bars:
                h = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2, h + 0.1,
                        f'{h:.1f}%', ha='center', fontsize=9, fontweight='bold')
        buf = self._chart_png(fig)
        slide.shapes.add_picture(buf, Inches(0.4), Inches(1.2), Inches(10.4), Inches(7.5))
        inter = tasas.get("interbancaria", "N/D")
        self._txt(slide, f"Tasa interbancaria (overnight): {inter}%",
                  m, Inches(10.2), w, Inches(0.5), size=12, color=GRAY, align=PP_ALIGN.CENTER)

    def _slide_top_cartera(self, prs, data):
        slide = self._blank_slide(prs, WHITE)
        m = Inches(0.3)
        w = SLIDE_SIZE - 2 * m
        self._txt(slide, "QUIEN CRECIO MAS?", m, Inches(0.3), w, Inches(0.7),
                  size=26, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        self._txt(slide, "Top entidades · mayor crecimiento de cartera de credito (i.a.)",
                  m, Inches(1.1), w, Inches(0.5), size=14, color=GRAY, align=PP_ALIGN.CENTER)
        top = data.get("sb", {}).get("rankings", {}).get("top_cartera_growth", [])[:5]
        if top:
            names = [e.get("entidad", f"E{i+1}") for i, e in enumerate(top)]
            values = [e.get("_growth_pct", 0) for e in top]
            colors = ['#1a1a2e', '#3182ce', '#63b3ed', '#90cdf4', '#bee3f8']
            fig, ax = plt.subplots(figsize=(9, 6))
            bars = ax.barh(names[::-1], values[::-1], color=colors[:len(top)], edgecolor='white')
            for bar, val in zip(bars, values[::-1]):
                ax.text(val + 0.2, bar.get_y() + bar.get_height()/2,
                        f'+{val:.1f}%', va='center', fontsize=11, fontweight='bold', color='#1a1a2e')
            ax.set_xlabel('Crecimiento i.a. (%)', fontsize=11)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            buf = self._chart_png(fig)
            slide.shapes.add_picture(buf, Inches(0.5), Inches(1.8), Inches(10.2), Inches(7.0))

    def _slide_top_npl(self, prs, data):
        slide = self._blank_slide(prs, WHITE)
        m = Inches(0.3)
        w = SLIDE_SIZE - 2 * m
        self._txt(slide, "MOROSIDAD POR ENTIDAD", m, Inches(0.3), w, Inches(0.7),
                  size=24, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        self._txt(slide, "Entidades con mayor incremento de cartera vencida (NPL, i.a. en pp)",
                  m, Inches(1.1), w, Inches(0.5), size=13, color=GRAY, align=PP_ALIGN.CENTER)
        top_npl = data.get("sb", {}).get("rankings", {}).get("top_npl_increase", [])[:5]
        y_pos = Inches(2.0)
        for i, entity in enumerate(top_npl):
            name = entity.get("entidad", f"Entidad {i+1}")
            val = entity.get("_growth_pct", 0)
            color = RED if val > 1 else YELLOW if val > 0.5 else GREEN
            indicator = slide.shapes.add_shape(1, Inches(0.5), y_pos, Inches(0.4), Inches(0.7))
            indicator.fill.solid()
            indicator.fill.fore_color.rgb = color
            indicator.line.color.rgb = color
            self._txt(slide, name, Inches(1.1), y_pos + Inches(0.05),
                      Inches(7.0), Inches(0.65), size=20, color=DARK_BLUE)
            self._txt(slide, f"+{val:.1f} pp", Inches(8.5), y_pos + Inches(0.05),
                      Inches(2.0), Inches(0.65), size=22, bold=True, color=color)
            y_pos += Inches(1.3)
        # Legend
        for x_off, color, label in [(Inches(0.5), GREEN, "Bajo (<0.5pp)"),
                                     (Inches(3.5), YELLOW, "Moderado"),
                                     (Inches(6.5), RED, "Elevado (>1pp)")]:
            dot = slide.shapes.add_shape(1, x_off, Inches(9.5), Inches(0.25), Inches(0.25))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.color.rgb = color
            self._txt(slide, label, x_off + Inches(0.35), Inches(9.45),
                      Inches(2.5), Inches(0.4), size=11, color=GRAY)

    def _slide_perspectivas(self, prs, data):
        slide = self._blank_slide(prs, DARK_BLUE)
        m = Inches(0.5)
        w = SLIDE_SIZE - 2 * m
        self._txt(slide, "PERSPECTIVAS", m, Inches(0.8), w, Inches(0.9),
                  size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        self._txt(slide, "Que vigilar el proximo mes?",
                  m, Inches(1.9), w, Inches(0.6), size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
        bcrd = data.get("bcrd", {})
        inf = bcrd.get("inflacion", {}).get("var_interanual") or 0
        imae = bcrd.get("imae", {}).get("var_interanual") or 0
        tpm = bcrd.get("tpm", {}).get("value") or 0
        points = [
            f"Inflacion en {inf:.1f}% i.a. — monitorear trayectoria vs. meta BCRD",
            f"IMAE +{imae:.1f}% — sostenibilidad del credito comercial?",
            f"TPM en {tpm:.2f}% — senales de cambio de postura monetaria",
            "Calidad de cartera ante expansion de credito de consumo",
        ]
        y = Inches(3.0)
        for point in points:
            self._txt(slide, f"• {point}", m, y, w, Inches(1.1), size=16, color=WHITE)
            y += Inches(1.5)

    def _slide_cta(self, prs, period):
        slide = self._blank_slide(prs, DARK_BLUE)
        m = Inches(0.5)
        w = SLIDE_SIZE - 2 * m
        self._txt(slide, "Que indicador te\ninteresa profundizar?",
                  m, Inches(2.5), w, Inches(2.5), size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        self._txt(slide, "Comenta abajo",
                  m, Inches(5.5), w, Inches(1.0), size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
        self._txt(slide, "#BancaDominicana #BCRD #FinanzasDominicanas #Macroeconomia",
                  m, Inches(9.8), w, Inches(0.7), size=12, color=GRAY, align=PP_ALIGN.CENTER)

    def generate(self, data: dict, period: str = None) -> str:
        if not period:
            period = date.today().strftime("%B %Y")
        prs = self._new_prs()
        self._slide_cover(prs, period)
        self._slide_kpis(prs, data)
        self._slide_macro_chart(prs, data)
        self._slide_tasas(prs, data)
        self._slide_top_cartera(prs, data)
        self._slide_top_npl(prs, data)
        self._slide_perspectivas(prs, data)
        self._slide_cta(prs, period)
        filename = f"{date.today()}-carousel.pptx"
        output_path = os.path.join(self.output_dir, filename)
        prs.save(output_path)
        return output_path
